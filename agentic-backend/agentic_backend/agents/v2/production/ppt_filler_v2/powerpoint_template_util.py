import logging
import re

from pptx import Presentation

logger = logging.getLogger(__name__)


def fill_slide_from_structured_response(ppt_path, structured_response, output_path):
    prs = Presentation(ppt_path)
    pattern = re.compile(r"\{([^}]+)\}")
    slide_numbers = [3, 4, 5]
    responses = structured_response.values()

    for slide_number, response in zip(slide_numbers, responses):
        slide = prs.slides[slide_number - 1]  # Convert to 0-based index

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            for paragraph in shape.text_frame.paragraphs:  # type: ignore
                # Merge all runs to find complete placeholders (handles split placeholders)
                full_text = "".join(run.text for run in paragraph.runs)

                # Find all placeholders and their positions in the merged text
                placeholder_replacements = {}
                for match in pattern.finditer(full_text):
                    key = match.group(1)
                    if key in response:
                        placeholder_replacements[match.start()] = {
                            "end": match.end(),
                            "placeholder": match.group(0),
                            "value": str(response[key]),
                        }

                if not placeholder_replacements:
                    continue

                # Sort placeholders by position to handle multiple placeholders correctly
                sorted_placeholders = sorted(
                    placeholder_replacements.items(), key=lambda x: x[0]
                )

                # Process each run and replace placeholder text
                char_position = 0
                for run in paragraph.runs:
                    run_start = char_position
                    run_end = char_position + len(run.text)
                    new_run_text = run.text
                    offset = 0  # Track text length changes within this run

                    for placeholder_start, replacement in sorted_placeholders:
                        placeholder_end = replacement["end"]

                        # Case 1: Placeholder starts in this run
                        if run_start <= placeholder_start < run_end:
                            local_start = placeholder_start - run_start + offset

                            if placeholder_end <= run_end:
                                # Placeholder is completely within this run
                                placeholder_len = len(replacement["placeholder"])
                                new_run_text = (
                                    new_run_text[:local_start]
                                    + replacement["value"]
                                    + new_run_text[local_start + placeholder_len :]
                                )
                                offset += len(replacement["value"]) - placeholder_len
                            else:
                                # Placeholder spans multiple runs - replace the start
                                new_run_text = (
                                    new_run_text[:local_start] + replacement["value"]
                                )

                        # Case 2: This run is in the middle/end of a placeholder
                        elif placeholder_start < run_start < placeholder_end:
                            if placeholder_end >= run_end:
                                # Entire run is part of the placeholder - empty it
                                new_run_text = ""
                            else:
                                # Placeholder ends in this run - keep text after placeholder
                                local_end = placeholder_end - run_start
                                new_run_text = new_run_text[local_end:]

                    run.text = new_run_text
                    char_position = run_end

    prs.save(output_path)
    return output_path
