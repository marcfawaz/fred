# Copyright Thales 2026
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""
Sample v2 agent: Slide Maker

Demonstrates the core v2 authoring patterns in a single file:

  1. ctx.config()             — read a typed configurable field
  2. ctx.read_resource()      — load a binary file from the agent config workspace
  3. ctx.extract_structured() — generate structured content with the agent LLM
  4. ctx.publish_bytes()      — publish a binary file and get a download link
  5. ctx.link()               — return the artifact as a download button in one line
  6. ui_field()               — declare a field once (no separate FieldSpec tuple)

The agent asks the user for a topic, generates a 5-sentence slide with the LLM,
fills a configurable .pptx template, and returns a clickable download link.

Compare with agents/v1/samples/slide_maker — same outcome, ~3× less code,
no LangGraph, no StateGraph, no manual state management.

Setup:
  Upload simple_template.pptx (included in this folder) to the agent config
  workspace via the agent edit drawer in the UI.
"""

from __future__ import annotations

import logging
import tempfile
from datetime import UTC, datetime
from pathlib import Path
from typing import Any, cast  # cast kept for pptx placeholder narrowing

from pptx import Presentation
from pptx.util import Pt
from pydantic import BaseModel, Field

from agentic_backend.core.agents.v2.authoring import (
    ReActAgent,
    ResourceNotFoundError,
    ToolContext,
    ToolOutput,
    UIHints,
    prompt_md,
    tool,
    ui_field,
)

logger = logging.getLogger(__name__)

_SLIDE_CONTENT_PROMPT = (
    "You are a concise summarization expert. "
    "Given a request or topic, produce a short slide title and a 5-sentence body paragraph "
    "suitable for a single presentation slide. No bullet points."
)


class SlideContent(BaseModel):
    """Structured slide content produced by the LLM."""

    title: str = Field(description="A short slide title (one line).")
    body: str = Field(description="5-sentence paragraph for the slide body.")


@tool(
    tool_ref="sample.slide_maker.generate",
    description=(
        "Generate a PowerPoint slide based on user instructions and return a download link. "
        "Reads the configured .pptx template, fills it with LLM-generated content, "
        "and publishes the result."
    ),
    success_message="Slide generated.",
)
async def generate_slide(ctx: ToolContext, instructions: str) -> ToolOutput:
    logger.info(f"generate_slide called with instructions: {instructions!r}")

    # 1. Generate structured slide content with the agent LLM
    logger.info("Extracting structured content...")
    slide_content = await ctx.extract_structured(
        SlideContent,
        prompt=_SLIDE_CONTENT_PROMPT,
        text=f"Request: {instructions}",
    )
    logger.info(f"Content extracted. Title: {slide_content.title!r}")

    # 2. Read the .pptx template from the agent config workspace
    template_key = ctx.config("ppt_template_key", default="simple_template.pptx")
    logger.info(f"Reading template resource: {template_key}")
    try:
        resource = await ctx.read_resource(template_key)
    except ResourceNotFoundError:
        logger.error(f"Template '{template_key}' not found.")
        return ctx.error(
            f"Template '{template_key}' not found. "
            "Upload it to the agent config workspace first."
        )

    # 3. Fill the template
    template_path: Path | None = None
    output_path: Path | None = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
            tmp.write(resource.content_bytes)
            template_path = Path(tmp.name)

        prs = Presentation(str(template_path))
        placeholder = cast(Any, prs.slides[0].placeholders[1])  # content placeholder
        if placeholder.has_text_frame:
            tf = placeholder.text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = f"{slide_content.title}\n\n{slide_content.body}"
            p.font.size = Pt(ctx.config("font_size_pt", default=14))

        with tempfile.NamedTemporaryFile(
            delete=False, suffix=".pptx", prefix="slide_"
        ) as out:
            prs.save(out.name)
            output_path = Path(out.name)

        # 4. Publish and return the download link
        timestamp = datetime.now(tz=UTC).strftime("%Y%m%d_%H%M%S")
        logger.info("Publishing generated slide artifact...")
        artifact = await ctx.publish_bytes(
            file_name=f"slide_{timestamp}.pptx",
            content=output_path.read_bytes(),
            content_type=(
                "application/vnd.openxmlformats-officedocument"
                ".presentationml.presentation"
            ),
            title=f"Slide: {slide_content.title}",
        )
        logger.info(f"Slide published. URL: {artifact}")
        return ctx.link(artifact, text=f"Slide '{slide_content.title}' generated.")

    finally:
        if template_path:
            template_path.unlink(missing_ok=True)
        if output_path:
            output_path.unlink(missing_ok=True)


class Definition(ReActAgent):
    """
    Sample v2 slide-maker agent.

    Generates a PowerPoint slide from any user topic and returns a download link.
    Upload simple_template.pptx to the agent config workspace before use.
    """

    agent_id: str = "sample.slide_maker.v2"
    role: str = "Slide Maker V2"
    description: str = (
        "Sample v2 agent: generates a PowerPoint slide from any topic "
        "and returns a download link. Demonstrates ctx.read_resource, "
        "ctx.extract_structured, and ctx.publish_bytes."
    )
    tags: tuple[str, ...] = ("sample", "academy", "powerpoint", "v2")
    tools = (generate_slide,)

    system_prompt_template: str = ui_field(
        prompt_md(
            package="agentic_backend.agents.v2.samples.slide_maker",
            file_name="slide_maker_system_prompt.md",
        ),
        title="System Prompt",
        description="Instructions for the slide generation assistant.",
        ui_type="prompt",
        required=True,
        ui=UIHints(group="Prompts", multiline=True, markdown=True),
    )

    ppt_template_key: str = ui_field(
        "simple_template.pptx",
        title="PowerPoint Template Key",
        description="Filename of the .pptx template stored in the agent config workspace.",
        ui_type="text",
        ui=UIHints(group="PowerPoint"),
    )
    font_size_pt: int = ui_field(
        14,
        title="Font Size (pt)",
        description="Font size applied to the generated slide text.",
        ui_type="number",
        ui=UIHints(group="PowerPoint"),
    )
    # `fields` is auto-derived from the ui_field() declarations above
