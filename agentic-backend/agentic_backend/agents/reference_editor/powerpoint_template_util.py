# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import asyncio
import logging
import re
from io import BytesIO

from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt

from agentic_backend.agents.reference_editor.image_search_util import (
    get_image_for_technology,
)

logger = logging.getLogger(__name__)


async def fill_slide_from_structured_response_async(
    ppt_path,
    structured_responses,
    output_path,
    vector_search_client=None,
    kf_base_client=None,
    search_options: dict | None = None,
):
    prs = Presentation(ppt_path)
    pattern = re.compile(r"\{([^}]+)\}")

    # Flatten the nested structure into a single dictionary
    flattened_data = {}
    for section_name, section_data in structured_responses.items():
        if isinstance(section_data, dict):
            # Add all nested key-value pairs to the flattened dictionary
            for key, value in section_data.items():
                flattened_data[key] = value
        else:
            # If it's not a dict, keep it as is
            flattened_data[section_name] = section_data

    logger.info(f"Flattened data keys: {list(flattened_data.keys())}")

    # Pre-fetch images for listeTechnologies if present
    tech_images = []
    if (
        "listeTechnologies" in flattened_data
        and vector_search_client
        and kf_base_client
    ):
        technologies_text = flattened_data["listeTechnologies"]
        logger.info(f"Detected listeTechnologies field with value: {technologies_text}")
        tech_images = await parse_technologies_and_fetch_images(
            technologies_text,
            vector_search_client,
            kf_base_client,
            search_options=search_options,
        )

    # Pre-fetch image for nomSociete if present
    societe_image = None
    if "nomSociete" in flattened_data and vector_search_client and kf_base_client:
        nom_societe = flattened_data["nomSociete"]
        logger.info(f"Detected nomSociete field with value: {nom_societe}")
        societe_image = await get_image_for_technology(
            nom_societe,
            vector_search_client,
            kf_base_client,
            search_options=search_options,
        )
        if societe_image:
            logger.info(f"Successfully fetched image for nomSociete: {nom_societe}")
        else:
            logger.info(f"No image found for nomSociete: {nom_societe}")

    # Process only the first slide (slide 0) for single-slide templates
    slide = prs.slides[0]

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:  # type: ignore
            # Merge all runs to find complete placeholders (handles split placeholders)
            full_text = "".join(run.text for run in paragraph.runs)

            # Find all placeholders and their positions in the merged text
            placeholder_replacements = {}
            has_liste_technologies = False
            has_nom_societe = False
            for match in pattern.finditer(full_text):
                key = match.group(1)
                if key in flattened_data:
                    # Check if this is the listeTechnologies placeholder
                    if key == "listeTechnologies" and tech_images:
                        has_liste_technologies = True
                        # We'll handle this separately with images
                        placeholder_replacements[match.start()] = {
                            "end": match.end(),
                            "placeholder": match.group(0),
                            "value": "",  # Clear the text, we'll add images
                        }
                    elif key == "nomSociete" and societe_image:
                        has_nom_societe = True
                        # Keep the text, we'll add the image next to it
                        placeholder_replacements[match.start()] = {
                            "end": match.end(),
                            "placeholder": match.group(0),
                            "value": str(flattened_data[key]),
                        }
                    else:
                        placeholder_replacements[match.start()] = {
                            "end": match.end(),
                            "placeholder": match.group(0),
                            "value": str(flattened_data[key]),
                        }
                else:
                    logger.warning(f"Placeholder '{key}' not found in flattened data")

            if not placeholder_replacements:
                continue

            # Sort placeholders by position to handle multiple placeholders correctly
            sorted_placeholders = sorted(
                placeholder_replacements.items(), key=lambda x: x[0]
            )

            # Process each run and replace placeholder text
            char_position = 0
            for run in paragraph.runs:
                run_start = char_position
                run_end = char_position + len(run.text)
                new_run_text = run.text
                offset = 0  # Track text length changes within this run

                for placeholder_start, replacement in sorted_placeholders:
                    placeholder_end = replacement["end"]

                    # Case 1: Placeholder starts in this run
                    if run_start <= placeholder_start < run_end:
                        local_start = placeholder_start - run_start + offset

                        if placeholder_end <= run_end:
                            # Placeholder is completely within this run
                            placeholder_len = len(replacement["placeholder"])
                            new_run_text = (
                                new_run_text[:local_start]
                                + replacement["value"]
                                + new_run_text[local_start + placeholder_len :]
                            )
                            offset += len(replacement["value"]) - placeholder_len
                        else:
                            # Placeholder spans multiple runs - replace the start
                            new_run_text = (
                                new_run_text[:local_start] + replacement["value"]
                            )

                    # Case 2: This run is in the middle/end of a placeholder
                    elif placeholder_start < run_start < placeholder_end:
                        if placeholder_end >= run_end:
                            # Entire run is part of the placeholder - empty it
                            new_run_text = ""
                        else:
                            # Placeholder ends in this run - keep text after placeholder
                            local_end = placeholder_end - run_start
                            new_run_text = new_run_text[local_end:]

                run.text = new_run_text
                char_position = run_end

            # After text replacement, if this shape contained listeTechnologies, add images
            if has_liste_technologies and tech_images:
                logger.info(
                    f"Adding {len(tech_images)} technology images to PowerPoint shape"
                )
                add_images_to_pptx_shape(shape, tech_images, slide)

            # After text replacement, if this shape contained nomSociete, add image next to it
            if has_nom_societe and societe_image:
                logger.info("Adding nomSociete image to PowerPoint shape")
                add_societe_image_to_pptx_shape(
                    shape, societe_image, flattened_data.get("nomSociete", ""), slide
                )

    prs.save(output_path)
    return output_path


def fill_slide_from_structured_response(
    ppt_path,
    structured_responses,
    output_path,
    vector_search_client=None,
    kf_base_client=None,
    search_options: dict | None = None,
):
    return asyncio.run(
        fill_slide_from_structured_response_async(
            ppt_path,
            structured_responses,
            output_path,
            vector_search_client,
            kf_base_client,
            search_options=search_options,
        )
    )


async def fill_word_from_structured_response_async(
    docx_path,
    structured_responses,
    output_path,
    vector_search_client=None,
    kf_base_client=None,
    search_options: dict | None = None,
):
    doc = Document(docx_path)
    pattern = re.compile(r"\{([^}]+)\}")

    logger.info("=== STARTING WORD TEMPLATE FILL ===")
    logger.info(f"Template path: {docx_path}")
    logger.info(f"Output path: {output_path}")

    # Log template content before processing
    logger.info("=== TEMPLATE CONTENT BEFORE PROCESSING ===")
    total_placeholders_in_template = 0
    for i, para in enumerate(doc.paragraphs):
        if para.text.strip():
            logger.info(f"Template Para {i}: {para.text}")
            placeholders = pattern.findall(para.text)
            if placeholders:
                total_placeholders_in_template += len(placeholders)
                logger.info(f"  Placeholders found: {placeholders}")

    for i, table in enumerate(doc.tables):
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                for para in cell.paragraphs:
                    if para.text.strip():
                        logger.info(
                            f"Template Table {i}[{row_idx},{col_idx}]: {para.text}"
                        )
                        placeholders = pattern.findall(para.text)
                        if placeholders:
                            total_placeholders_in_template += len(placeholders)
                            logger.info(f"  Placeholders found: {placeholders}")

    logger.info(f"Total placeholders in template: {total_placeholders_in_template}")

    # Flatten the nested structure into a single dictionary
    flattened_data = {}
    for section_name, section_data in structured_responses.items():
        if isinstance(section_data, dict):
            # Add all nested key-value pairs to the flattened dictionary
            for key, value in section_data.items():
                flattened_data[key] = value
        else:
            # If it's not a dict, keep it as is
            flattened_data[section_name] = section_data

    logger.info("=== DATA FOR REPLACEMENT ===")
    logger.info(f"Flattened data keys: {list(flattened_data.keys())}")
    logger.info(f"Flattened data: {flattened_data}")

    # Pre-fetch images for listeTechnologies if present
    tech_images = []
    if (
        "listeTechnologies" in flattened_data
        and vector_search_client
        and kf_base_client
    ):
        technologies_text = flattened_data["listeTechnologies"]
        logger.info(f"Detected listeTechnologies field with value: {technologies_text}")
        tech_images = await parse_technologies_and_fetch_images(
            technologies_text,
            vector_search_client,
            kf_base_client,
            search_options=search_options,
        )

    # Pre-fetch image for nomSociete if present
    societe_image = None
    if "nomSociete" in flattened_data and vector_search_client and kf_base_client:
        nom_societe = flattened_data["nomSociete"]
        logger.info(f"Detected nomSociete field with value: {nom_societe}")
        societe_image = await get_image_for_technology(
            nom_societe,
            vector_search_client,
            kf_base_client,
            search_options=search_options,
        )
        if societe_image:
            logger.info(f"Successfully fetched image for nomSociete: {nom_societe}")
        else:
            logger.info(f"No image found for nomSociete: {nom_societe}")

    # Helper function to process paragraphs with formatting preservation
    def process_paragraph(paragraph):
        # Merge all runs to find complete placeholders (handles split placeholders)
        full_text = "".join(run.text for run in paragraph.runs)

        if not full_text:
            return

        logger.debug(f"Processing paragraph text: '{full_text}'")

        # Find all placeholders and their positions in the merged text
        matches = list(pattern.finditer(full_text))
        if not matches:
            return

        logger.info(f"Found {len(matches)} placeholders in paragraph")

        # Build a list of placeholder replacements
        placeholder_replacements = {}
        has_liste_technologies = False
        has_nom_societe = False
        for match in matches:
            key = match.group(1)
            logger.info(f"Processing placeholder: {{{key}}}")
            if key in flattened_data:
                # Check if this is the listeTechnologies placeholder
                if key == "listeTechnologies" and tech_images:
                    has_liste_technologies = True
                    # We'll handle this separately with images
                    placeholder_replacements[match.start()] = {
                        "end": match.end(),
                        "placeholder": match.group(0),
                        "value": "",  # Clear the text, we'll add images
                    }
                elif key == "nomSociete" and societe_image:
                    has_nom_societe = True
                    # Keep the text, we'll add the image next to it
                    value = str(flattened_data[key])
                    logger.info(f"Replacing {{{key}}} with: {value} (will add image)")
                    placeholder_replacements[match.start()] = {
                        "end": match.end(),
                        "placeholder": match.group(0),
                        "value": value,
                    }
                else:
                    value = str(flattened_data[key])
                    logger.info(f"Replacing {{{key}}} with: {value}")
                    placeholder_replacements[match.start()] = {
                        "end": match.end(),
                        "placeholder": match.group(0),
                        "value": value,
                    }
            else:
                logger.warning(f"Placeholder '{key}' not found in flattened data")

        if not placeholder_replacements:
            return

        # Sort placeholders by position
        sorted_placeholders = sorted(
            placeholder_replacements.items(), key=lambda x: x[0]
        )

        # Process each run and replace placeholder text while preserving formatting
        char_position = 0
        for run in paragraph.runs:
            run_start = char_position
            run_end = char_position + len(run.text)
            new_run_text = run.text
            offset = 0  # Track text length changes within this run

            for placeholder_start, replacement in sorted_placeholders:
                placeholder_end = replacement["end"]

                # Case 1: Placeholder starts in this run
                if run_start <= placeholder_start < run_end:
                    local_start = placeholder_start - run_start + offset

                    if placeholder_end <= run_end:
                        # Placeholder is completely within this run
                        placeholder_len = len(replacement["placeholder"])
                        new_run_text = (
                            new_run_text[:local_start]
                            + replacement["value"]
                            + new_run_text[local_start + placeholder_len :]
                        )
                        offset += len(replacement["value"]) - placeholder_len
                    else:
                        # Placeholder spans multiple runs - replace the start
                        new_run_text = new_run_text[:local_start] + replacement["value"]
                        offset = len(new_run_text) - local_start

                # Case 2: This run is in the middle/end of a placeholder
                elif placeholder_start < run_start < placeholder_end:
                    if placeholder_end >= run_end:
                        # Entire run is part of the placeholder - empty it
                        new_run_text = ""
                        offset = -len(run.text)
                    else:
                        # Placeholder ends in this run - keep text after placeholder
                        local_end = placeholder_end - run_start
                        new_run_text = new_run_text[local_end:]
                        offset = -local_end

            run.text = new_run_text
            char_position = run_end

        # After text replacement, if this paragraph contained listeTechnologies, add images
        if has_liste_technologies and tech_images:
            logger.info(
                f"Adding {len(tech_images)} technology images to Word paragraph"
            )
            add_images_to_word_paragraph(paragraph, tech_images)

        # After text replacement, if this paragraph contained nomSociete, add image next to it
        if has_nom_societe and societe_image:
            logger.info("Adding nomSociete image to Word paragraph")
            add_societe_image_to_word_paragraph(
                paragraph, societe_image, flattened_data.get("nomSociete", "")
            )

        logger.info(f"Final paragraph text after replacement: '{paragraph.text}'")

    # Process all paragraphs in the document body
    for paragraph in doc.paragraphs:
        process_paragraph(paragraph)

    # Process tables in the document
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    process_paragraph(paragraph)

    # Process headers
    for section in doc.sections:
        # Process header (first, even, default)
        header = section.header
        for paragraph in header.paragraphs:
            process_paragraph(paragraph)
        # Process header tables
        for table in header.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        process_paragraph(paragraph)

    # Process footers
    for section in doc.sections:
        # Process footer (first, even, default)
        footer = section.footer
        for paragraph in footer.paragraphs:
            process_paragraph(paragraph)
        # Process footer tables
        for table in footer.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        process_paragraph(paragraph)

    # Process textboxes (VML and DrawingML) - CRITICAL for templates with text boxes
    # Many Word templates use textboxes for layout, and python-docx doesn't handle them well
    logger.info("=== PROCESSING TEXTBOXES ===")

    textbox_count = 0
    replacements_made = 0
    textbox_with_liste_tech = None  # Store textbox element containing listeTechnologies
    textbox_with_nom_societe = None  # Store textbox element containing nomSociete

    # Process all textbox content elements (both VML and DrawingML formats)
    for txbxContent in doc.element.body.iter():
        # Check if this is a textbox content element
        tag_name = (
            txbxContent.tag.split("}")[-1]
            if "}" in txbxContent.tag
            else txbxContent.tag
        )
        if tag_name != "txbxContent":
            continue

        textbox_count += 1
        logger.info(f"Processing textbox {textbox_count}")

        # Find all text elements in this textbox
        for t_elem in txbxContent.iter():
            t_tag = t_elem.tag.split("}")[-1] if "}" in t_elem.tag else t_elem.tag
            if t_tag != "t":
                continue

            if t_elem.text and ("{" in t_elem.text or "}" in t_elem.text):
                original_text = t_elem.text
                logger.info(
                    f"  Found text with potential placeholders: '{original_text}'"
                )

                # Find and replace all placeholders in this text element
                new_text = original_text
                for match in pattern.finditer(original_text):
                    key = match.group(1)
                    placeholder = match.group(0)

                    if key in flattened_data:
                        # Special case: listeTechnologies with images
                        if key == "listeTechnologies" and tech_images:
                            logger.info(
                                "    Found listeTechnologies, will add images separately"
                            )
                            new_text = new_text.replace(placeholder, "")
                            textbox_with_liste_tech = txbxContent
                            replacements_made += 1
                        elif key == "nomSociete" and societe_image:
                            # Replace with value, and mark for image addition
                            value = str(flattened_data[key])
                            logger.info(
                                f"    Replacing {placeholder} with: {value} (will add image)"
                            )
                            new_text = new_text.replace(placeholder, value)
                            textbox_with_nom_societe = txbxContent
                            replacements_made += 1
                        else:
                            value = str(flattened_data[key])
                            logger.info(f"    Replacing {placeholder} with: {value}")
                            new_text = new_text.replace(placeholder, value)
                            replacements_made += 1
                    else:
                        logger.warning(
                            f"    Placeholder '{key}' not found in flattened data"
                        )

                if new_text != original_text:
                    t_elem.text = new_text
                    logger.info(f"    Updated text: '{new_text}'")

    logger.info(
        f"Processed {textbox_count} textboxes, made {replacements_made} replacements"
    )

    # Add images near the textbox that contained listeTechnologies
    if textbox_with_liste_tech is not None and tech_images:
        logger.info(f"Adding {len(tech_images)} images near listeTechnologies textbox")
        _add_images_near_textbox(doc, textbox_with_liste_tech, tech_images)

    # Add image near the textbox that contained nomSociete
    if textbox_with_nom_societe is not None and societe_image:
        logger.info("Adding nomSociete image near textbox")
        _add_societe_image_near_textbox(
            doc,
            textbox_with_nom_societe,
            societe_image,
            flattened_data.get("nomSociete", ""),
        )

    doc.save(output_path)
    return output_path


def fill_word_from_structured_response(
    docx_path,
    structured_responses,
    output_path,
    vector_search_client=None,
    kf_base_client=None,
    search_options: dict | None = None,
):
    return asyncio.run(
        fill_word_from_structured_response_async(
            docx_path,
            structured_responses,
            output_path,
            vector_search_client,
            kf_base_client,
            search_options=search_options,
        )
    )


async def parse_technologies_and_fetch_images(
    technologies_text: str,
    vector_search_client,
    kf_base_client,
    *,
    search_options: dict | None = None,
) -> list[tuple[str, BytesIO | None]]:
    """
    Parse a comma-separated list of technologies and fetch their corresponding images.

    Args:
        technologies_text: Comma-separated list of technology names (e.g., "Nvidia, Apple, AWS, SharePoint")
        kf_client: Authenticated KfBaseClient instance

    Returns:
        List of tuples (technology_name, image_data), where image_data is BytesIO or None if not found
    """
    if not technologies_text or not isinstance(technologies_text, str):
        return []

    # Split by commas and clean each technology name
    technologies = [
        tech.strip() for tech in technologies_text.split(",") if tech.strip()
    ]

    logger.info(f"Parsing {len(technologies)} technologies: {technologies}")

    results = []
    for tech_name in technologies:
        logger.info(f"Fetching image for technology: {tech_name}")
        image_data = await get_image_for_technology(
            tech_name,
            vector_search_client,
            kf_base_client,
            search_options=search_options,
        )

        if image_data:
            logger.info(f"Successfully fetched image for: {tech_name}")
            results.append((tech_name, image_data))
        else:
            logger.warning(
                f"Could not fetch image for: {tech_name}, will use text instead"
            )
            results.append((tech_name, None))

    logger.info(
        f"Fetched {sum(1 for _, img in results if img is not None)} images out of {len(technologies)} technologies"
    )
    return results


def add_images_to_pptx_shape(shape, images: list[tuple[str, BytesIO | None]], slide):
    """
    Add images to a PowerPoint shape placeholder, replacing the text content.

    Args:
        shape: The PowerPoint shape containing the placeholder
        images: List of tuples (technology_name, image_data)
        slide: The slide object to add images to
    """
    if not images:
        return

    # Clear the text content from the shape
    if shape.has_text_frame:
        shape.text_frame.clear()

    # Calculate positioning for images
    # Get the shape's position and size
    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height

    # Calculate image dimensions (we'll arrange them horizontally)
    num_images = len([img for _, img in images if img is not None])
    if num_images == 0:
        # No images, just set text
        shape.text = ", ".join([name for name, _ in images])
        return

    # Image size and spacing
    image_width = width // num_images
    image_height = min(height, PptxInches(0.8))  # Max height of 0.8 inches

    # Add each image
    current_left = left
    for tech_name, image_data in images:
        if image_data is not None:
            try:
                # Reset the BytesIO pointer to the beginning
                image_data.seek(0)

                # Add the image to the slide
                slide.shapes.add_picture(
                    image_data,
                    current_left,
                    top + (height - image_height) // 2,  # Center vertically
                    width=image_width,
                    height=image_height,
                )
                current_left += image_width
                logger.info(f"Added image for technology: {tech_name}")
            except Exception as e:
                logger.error(f"Failed to add image for {tech_name}: {e}")


def add_societe_image_to_pptx_shape(
    shape, societe_image: BytesIO, nom_societe: str, slide
):
    """
    Add the company image next to the nomSociete text in a PowerPoint shape.
    The image is placed at the beginning of the shape area (like a logo before the company name).

    Args:
        shape: The PowerPoint shape containing the nomSociete text
        societe_image: BytesIO object containing the image data
        nom_societe: The name of the company (for logging)
        slide: The slide object to add the image to
    """
    if not societe_image:
        logger.warning(f"No societe_image provided for {nom_societe}")
        return

    try:
        # Make a copy of the image data to avoid issues with BytesIO pointer
        societe_image.seek(0)
        image_bytes = societe_image.read()
        image_copy = BytesIO(image_bytes)

        # Get the shape's position and size
        shape_left = shape.left
        shape_top = shape.top
        shape_height = shape.height

        # Use specified image dimensions: 69x56 points (343/5 x 280/5)
        image_width = PptxPt(69)
        image_height = PptxPt(56)

        # Position the image at the start of the shape (inside the shape area)
        image_left = shape_left
        image_top = shape_top + (shape_height - image_height) // 2

        # Ensure image_top is not negative
        if image_top < 0:
            image_top = shape_top

        logger.info(
            f"Adding company image for {nom_societe}: left={image_left}, top={image_top}, width={image_width}, height={image_height}"
        )
        logger.info(
            f"Shape info: left={shape_left}, top={shape_top}, height={shape_height}"
        )

        # Add the image to the slide
        picture = slide.shapes.add_picture(
            image_copy, image_left, image_top, width=image_width, height=image_height
        )
        logger.info(
            f"Successfully added company image for: {nom_societe}, picture shape id: {picture.shape_id}"
        )
    except Exception as e:
        logger.error(
            f"Failed to add company image for {nom_societe}: {e}", exc_info=True
        )


def add_images_to_word_paragraph(paragraph, images: list[tuple[str, BytesIO | None]]):
    """
    Add images to a Word paragraph, replacing the text content.

    Args:
        paragraph: The Word paragraph object
        images: List of tuples (technology_name, image_data)
    """
    if not images:
        return

    # Clear existing runs
    for run in paragraph.runs:
        run.text = ""

    # Add each image or text
    for i, (tech_name, image_data) in enumerate(images):
        if image_data is not None:
            try:
                # Reset the BytesIO pointer
                image_data.seek(0)

                # Add image to the paragraph
                run = paragraph.add_run()
                run.add_picture(image_data, width=Inches(0.5))

                # Add spacing between images
                if i < len(images) - 1:
                    paragraph.add_run("  ")

                logger.info(f"Added image for technology: {tech_name} to Word document")
            except Exception as e:
                logger.error(f"Failed to add image for {tech_name} to Word: {e}")
                # Fallback to text
                paragraph.add_run(tech_name)
                if i < len(images) - 1:
                    paragraph.add_run(", ")
        else:
            # No image found, use text
            paragraph.add_run(tech_name)
            if i < len(images) - 1:
                paragraph.add_run(", ")


def add_societe_image_to_word_paragraph(
    paragraph, societe_image: BytesIO, nom_societe: str
):
    """
    Add the company image next to the nomSociete text in a Word paragraph.

    Args:
        paragraph: The Word paragraph object
        societe_image: BytesIO object containing the image data
        nom_societe: The name of the company (for logging)
    """
    if not societe_image:
        return

    try:
        # Reset the BytesIO pointer
        societe_image.seek(0)

        # Add a space and then the image at the end of the paragraph
        # Use specified image dimensions: 69x56 points (343/5 x 280/5)
        run = paragraph.add_run("  ")  # Add spacing
        run = paragraph.add_run()
        run.add_picture(societe_image, width=Pt(69), height=Pt(56))

        logger.info(
            f"Added company image for: {nom_societe} to Word paragraph (69x56 pt)"
        )
    except Exception as e:
        logger.error(f"Failed to add company image for {nom_societe} to Word: {e}")


def _add_images_near_textbox(
    doc, textbox_element, images: list[tuple[str, BytesIO | None]]
):
    """
    Add images in a new paragraph right after the textbox containing {listeTechnologies}.
    """
    if not images:
        return

    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement

    # Find the parent paragraph in the document body that contains this textbox
    parent = textbox_element.getparent()
    containing_paragraph = None

    while parent is not None:
        tag_name = parent.tag.split("}")[-1] if "}" in parent.tag else parent.tag
        if tag_name == "p":
            grandparent = parent.getparent()
            if grandparent is not None:
                gp_tag = (
                    grandparent.tag.split("}")[-1]
                    if "}" in grandparent.tag
                    else grandparent.tag
                )
                if gp_tag == "body":
                    containing_paragraph = parent
                    break
        parent = parent.getparent()

    if containing_paragraph is None:
        logger.warning("Could not find containing paragraph, adding images at end")
        paragraph = doc.add_paragraph()
    else:
        body = containing_paragraph.getparent()
        para_index = list(body).index(containing_paragraph)
        logger.info(f"Inserting images after paragraph at index {para_index}")

        new_p = OxmlElement("w:p")
        body.insert(para_index + 1, new_p)

        paragraph = None
        for p in doc.paragraphs:
            if p._element is new_p:
                paragraph = p
                break

        if paragraph is None:
            logger.warning("Could not get Paragraph object, adding at end")
            paragraph = doc.add_paragraph()

    # Align to the right
    paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT

    # Filter to only keep images that have data (skip technologies without images)
    images_with_data = [(name, data) for name, data in images if data is not None]

    if not images_with_data:
        logger.info("No images with data, skipping")
        return

    # Add only images (no text fallback for missing images)
    for i, (tech_name, image_data) in enumerate(images_with_data):
        try:
            image_data.seek(0)
            run = paragraph.add_run()
            run.add_picture(image_data, width=Inches(0.5))
            if i < len(images_with_data) - 1:
                paragraph.add_run("  ")
            logger.info(f"Added image for: {tech_name}")
        except Exception as e:
            logger.error(f"Failed to add image for {tech_name}: {e}")


def _add_societe_image_near_textbox(
    doc, textbox_element, societe_image: BytesIO, nom_societe: str
):
    """
    Add the company image in a new paragraph right after the textbox containing {nomSociete}.
    """
    if not societe_image:
        return

    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml import OxmlElement

    # Find the parent paragraph in the document body that contains this textbox
    parent = textbox_element.getparent()
    containing_paragraph = None

    while parent is not None:
        tag_name = parent.tag.split("}")[-1] if "}" in parent.tag else parent.tag
        if tag_name == "p":
            grandparent = parent.getparent()
            if grandparent is not None:
                gp_tag = (
                    grandparent.tag.split("}")[-1]
                    if "}" in grandparent.tag
                    else grandparent.tag
                )
                if gp_tag == "body":
                    containing_paragraph = parent
                    break
        parent = parent.getparent()

    if containing_paragraph is None:
        logger.warning(
            "Could not find containing paragraph for nomSociete, adding image at end"
        )
        paragraph = doc.add_paragraph()
    else:
        body = containing_paragraph.getparent()
        para_index = list(body).index(containing_paragraph)
        logger.info(f"Inserting nomSociete image after paragraph at index {para_index}")

        new_p = OxmlElement("w:p")
        body.insert(para_index + 1, new_p)

        paragraph = None
        for p in doc.paragraphs:
            if p._element is new_p:
                paragraph = p
                break

        if paragraph is None:
            logger.warning(
                "Could not get Paragraph object for nomSociete, adding at end"
            )
            paragraph = doc.add_paragraph()

    # Align to the left (near the text)
    paragraph.alignment = WD_ALIGN_PARAGRAPH.LEFT

    try:
        societe_image.seek(0)
        run = paragraph.add_run()
        # Use specified image dimensions: 69x56 points (343/5 x 280/5)
        run.add_picture(societe_image, width=Pt(69), height=Pt(56))
        logger.info(f"Added company image for: {nom_societe} near textbox (69x56 pt)")
    except Exception as e:
        logger.error(f"Failed to add company image for {nom_societe} near textbox: {e}")


referenceSchema = {
    "type": "object",
    "required": ["informationsProjet", "contexte", "syntheseProjet"],
    "additionalProperties": False,
    "properties": {
        "informationsProjet": {
            "type": "object",
            "description": "Informations sur le projet",
            "required": [
                "nomSociete",
                "nomProjet",
                "dateProjet",
                "nombrePersonnes",
                "enjeuFinancier",
            ],
            "additionalProperties": False,
            "properties": {
                "nomSociete": {
                    "type": "string",
                    "description": "Le nom de la société",
                    "maxLength": 50,
                },
                "nomProjet": {
                    "type": "string",
                    "description": "Le nom du projet",
                    "maxLength": 50,
                },
                "dateProjet": {
                    "type": "string",
                    "description": "La date de début et de fin du projet",
                },
                "nombrePersonnes": {
                    "type": "string",
                    "description": "Le nombre de personnes dans la direction (uniquement dans la direction)",
                },
                "enjeuFinancier": {
                    "type": "string",
                    "description": "Les coûts financiers ou la rentabilité du projet exprimé en euros, juste un seul chiffre clé comme le CA (jamais en nombre de personnes sinon ne rien mettre)",
                    "maxLength": 100,
                },
            },
        },
        "contexte": {
            "type": "object",
            "description": "Informations sur le contexte et le client",
            "required": [
                "presentationClient",
                "presentationContexte",
                "listeTechnologies",
            ],
            "additionalProperties": False,
            "properties": {
                "presentationClient": {
                    "type": "string",
                    "description": "Courte présentation du client. Longueur maximale: 300 caractères (une à deux phrases).",
                    "maxLength": 300,
                },
                "presentationContexte": {
                    "type": "string",
                    "description": "Courte présentation du contexte du projet. Longueur maximale: 300 caractères (une à deux phrases).",
                    "maxLength": 300,
                },
                "listeTechnologies": {
                    "type": "string",
                    "description": "Listes des technologies utilisés lors de ce projet",
                },
            },
        },
        "syntheseProjet": {
            "type": "object",
            "description": "Synthèse struturée du projet",
            "required": [
                "enjeux",
                "activiteSolutions",
                "beneficeClients",
                "pointsForts",
            ],
            "additionalProperties": False,
            "properties": {
                "enjeux": {
                    "type": "string",
                    "description": "Court résumé des enjeux du projet. Longueur maximale: 300 caractères (une à deux phrases).",
                    "maxLength": 300,
                },
                "activiteSolutions": {
                    "type": "string",
                    "description": "Court résumé des activités et solutions du projet. Longueur maximale: 300 caractères (une à deux phrases).",
                    "maxLength": 300,
                },
                "beneficeClients": {
                    "type": "string",
                    "description": "Court résumé des bénéfices pour le client. Longueur maximale: 300 caractères (une à deux phrases).",
                },
                "pointsForts": {
                    "type": "string",
                    "description": "Court résumé des points forts du projet. Longueur maximale: 300 caractères (une à deux phrases).",
                },
            },
        },
    },
}
