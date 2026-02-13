# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

import logging
import os
import shutil
import subprocess  # nosec: controlled subprocess usage
import tempfile
import time
from pathlib import Path
from typing import Dict, List, Optional

from fred_core import VectorSearchHit
from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
from langchain_core.output_parsers.json import JsonOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langgraph.graph import END, START, MessagesState, StateGraph
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pydantic import BaseModel

from agentic_backend.application_context import get_default_chat_model
from agentic_backend.common.kf_vectorsearch_client import VectorSearchClient
from agentic_backend.common.kf_workspace_client import WorkspaceRetrievalError
from agentic_backend.common.rags_utils import (
    attach_sources_to_llm_response,
    ensure_ranks,
    format_sources_for_prompt,
    sort_hits,
)
from agentic_backend.core.agents.agent_flow import AgentFlow
from agentic_backend.core.agents.agent_spec import AgentTuning, FieldSpec, UIHints
from agentic_backend.core.agents.runtime_context import (
    RuntimeContext,
    get_document_library_tags_ids,
    get_document_uids,
    get_search_policy,
    get_vector_search_scopes,
)
from agentic_backend.core.chatbot.chat_schema import (
    LinkKind,
    LinkPart,
    MessagePart,
    TextPart,
)
from agentic_backend.core.runtime_source import expose_runtime_source

logger = logging.getLogger(__name__)

# -----------------------------
# Tuning (UI schema)
# -----------------------------
RAG_TUNING = AgentTuning(
    role="slide_maker",
    description="An agent that generates PowerPoint slides with LLM content and provides a download link.",
    fields=[
        FieldSpec(
            key="prompts.system",
            type="prompt",
            title="RAG System Prompt",
            description="Defines assistant behavior and citation style. Example: 'Always reply in French.'",
            required=True,
            default=(
                "You answer strictly based on the retrieved document chunks.\n"
                "Extract the project reference information\n"
                "Each field must be concise (3-5 short items or a brief sentence)."
            ),
            ui=UIHints(group="Prompts", multiline=True, markdown=True),
        ),
        # NEW FIELD 1: Structured Extraction Prompt Template
        FieldSpec(
            key="prompts.extraction_template",
            type="prompt",
            title="Structured Extraction Prompt",
            description="Template for extracting structured data from RAG context. MUST contain {format_instructions} and {context} placeholders.",
            required=True,
            default="""You are an assistant that reads project reference documents.
Extract the following information. Ensure the content for all description fields (6-10) is translated into concise French before being put into the structured format:

HEADER INFORMATION:
1. Project name
2. Start date (month and year if possible)
3. End date (month and year if possible)
4. Number of people involved
5. Budget in k€

BODY INFORMATION (Must be in French):
6. Client presentation and context
7. Project stakes
8. Activities and solutions
9. Benefits for the client
10. Project strengths
11. List of technologies used

No field should remain empty unless the information is completely missing.
Each field must be concise (3-5 short items or a brief sentence).

Return strictly in this structured format:
{format_instructions}

Context:
{context}""",
            ui=UIHints(group="Prompts", multiline=True),
        ),
        # NEW FIELD 2: Summary Prompt Template
        FieldSpec(
            key="prompts.summary_template",
            type="prompt",
            title="Final Summary Prompt",
            description="Template used to generate a chat summary of the structured data. MUST contain {structured_data} placeholder.",
            required=True,
            default="""You are an assistant that summarizes project reference information for a chat user.
Given the extracted structured data below, produce a concise, readable summary in 3-5 short sentences.You must always respond in the same language as the source documents you extract information from.
If multiple languages are detected, use the predominant one. If non are detected fall back to French.
Do not repeat fields unnecessarily, focus on key points like project name, duration, team size, budget, and main activities.
Structured data:
{structured_data}""",
            ui=UIHints(group="Prompts", multiline=True),
        ),
        FieldSpec(
            key="rag.top_k",
            type="integer",
            title="Top-K Documents",
            description="Number of chunks to retrieve per question.",
            required=False,
            default=6,
            ui=UIHints(group="Retrieval"),
        ),
        FieldSpec(
            key="prompts.include_chat_context",
            type="boolean",
            title="Append Chat Context to System Prompt",
            description="If true, append the runtime chat context text after the system prompt.",
            required=False,
            default=True,
            ui=UIHints(group="Prompts"),
        ),
        FieldSpec(
            key="ppt.template_path",
            type="text",
            title="PowerPoint Template Path",
            description="Filesystem path to the .pptx template used to generate the fiche.",
            required=False,
            default="template_fiche_ref_projet.pptx",
            ui=UIHints(group="PowerPoint"),
        ),
    ],
)


@expose_runtime_source("agent.Sloan")
class Sloan(AgentFlow):
    """
    RAG agent for project reference extraction:
    - Retrieves context chunks via VectorSearchClient
    - Parses structured data fields
    - Fills a PowerPoint template and exports both .pptx and .pdf
    - Returns download + preview links
    """

    tuning = RAG_TUNING

    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        # Initialize placeholders for prompts
        self.parser: Optional[JsonOutputParser] = None
        self.prompt_template: Optional[ChatPromptTemplate] = None
        self.summary_prompt_template: Optional[ChatPromptTemplate] = None
        self.format_instructions: str = ""

    async def async_init(self, runtime_context: RuntimeContext):
        await super().async_init(runtime_context)
        self.model = get_default_chat_model()
        self.search_client = VectorSearchClient(agent=self)
        self._graph = self._build_graph()

        # Structured output schema using a Pydantic model (JSON output)
        class ProjectReferenceSchema(BaseModel):
            project_name: str
            start_date: str
            end_date: str
            num_people: str
            budget_k_eur: str
            client_presentation_and_context: str
            stakes: str
            activities_and_solutions: str
            client_benefits: str
            strengths: str
            tech_list: str

        self.parser = JsonOutputParser(pydantic_object=ProjectReferenceSchema)
        self.format_instructions = self.parser.get_format_instructions()

        # --- Use Tuned Prompts ---
        extraction_template_text = self.get_tuned_text("prompts.extraction_template")
        summary_template_text = self.get_tuned_text("prompts.summary_template")

        if not extraction_template_text or not summary_template_text:
            raise RuntimeError(
                "Missing required extraction or summary prompt template in tuning."
            )

        # Extraction Prompt (for structured data generation)
        self.prompt_template = ChatPromptTemplate.from_template(
            extraction_template_text
        )

        # Summary Prompt (for final chat output summary)
        self.summary_prompt_template = ChatPromptTemplate.from_template(
            summary_template_text
        )
        # -------------------------
        logger.info("Sloan initialized with structured output parsing enabled.")

    # -----------------------------
    # Graph
    # -----------------------------
    def _build_graph(self) -> StateGraph:
        builder = StateGraph(MessagesState)
        builder.add_node("reasoner", self._run_reasoning_step)
        builder.add_edge(START, "reasoner")
        builder.add_edge("reasoner", END)
        return builder

    # -----------------------------
    # Helpers
    # -----------------------------

    def _convert_pptx_to_pdf(self, pptx_path: Path) -> Optional[Path]:
        """Convert PPTX to PDF using headless LibreOffice with font embedding."""
        pdf_path = pptx_path.with_suffix(".pdf")
        try:
            soffice_path = shutil.which("soffice")
            if not soffice_path:
                raise FileNotFoundError(
                    "LibreOffice executable 'soffice' not found in PATH. "
                    "Please ensure LibreOffice is installed and 'soffice' is in your PATH."
                )

            subprocess.run(
                [
                    soffice_path,
                    "--headless",
                    "--nologo",
                    "--nofirststartwizard",
                    "--convert-to",
                    "pdf:writer_pdf_Export:EmbedStandardFonts=True,SelectPdfVersion=1",
                    "--outdir",
                    str(pptx_path.parent),
                    str(pptx_path),
                ],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
            )  # nosec: inputs are fully controlled and shell=False

            if pdf_path.exists():
                logger.info(
                    "LibreOffice successfully converted PPTX to PDF with embedded fonts: %s",
                    pdf_path,
                )
                return pdf_path
            else:
                logger.warning(
                    "LibreOffice conversion completed but PDF not found at %s", pdf_path
                )
                return None
        except subprocess.CalledProcessError as e:
            logger.error(
                "LibreOffice PDF conversion failed: %s",
                e.stderr.decode(errors="ignore"),
            )
            return None
        except FileNotFoundError:
            logger.error("LibreOffice (soffice) is not installed or not in PATH.")
            return None

    def _system_prompt(self) -> str:
        """Retrieve tuned system prompt."""
        sys_tpl = self.get_tuned_text("prompts.system")
        if not sys_tpl:
            logger.warning("Sloan: no tuned system prompt found.")
            raise RuntimeError("Missing system prompt.")
        return self.render(sys_tpl)

    # -----------------------------
    # PowerPoint filling
    # -----------------------------
    async def _fill_ppt_template(
        self, output_data: Dict[str, str], project_name: str
    ) -> Path:
        """Fill PowerPoint template and style header, with fallback text for missing info."""

        template_key = (
            self.get_tuned_text("ppt.template_path") or "template_fiche_ref_projet.pptx"
        )

        # If the key already ends with .pptx, don't append another suffix
        suffix = "" if template_key.lower().endswith(".pptx") else ".pptx"

        template_path = await self.fetch_config_blob_to_tempfile(
            template_key, suffix=suffix
        )

        if not os.path.exists(template_path):
            raise FileNotFoundError(f"PowerPoint template not found: {template_path}")

        # Helper to normalize missing fields
        def safe_get(key: str, suffix: str = "") -> str:
            value = str(output_data.get(key, "") or "").strip()
            if not value:
                value = "N/A"
            return f"{value}{suffix}" if suffix else value

        with tempfile.NamedTemporaryFile(
            delete=False, suffix=".pptx", prefix=f"fiche_{project_name}_"
        ) as out:
            output_path = Path(out.name)

        prs = Presentation(str(template_path))
        slide = prs.slides[0]

        # --- BODY placeholders ---
        mapping = {
            "client_presentation_and_context": 10,
            "stakes": 11,
            "activities_and_solutions": 12,
            "client_benefits": 13,
            "strengths": 14,
            "tech_list": 15,
        }

        for key, ph_idx in mapping.items():
            try:
                placeholder = slide.placeholders[ph_idx]
                if not getattr(placeholder, "has_text_frame", False):
                    continue
                textbox = placeholder.text_frame  # type: ignore
                textbox.clear()
                p = textbox.add_paragraph()
                p.text = safe_get(key)
                p.font.size = Pt(10)
            except Exception as e:
                logger.warning(
                    "Error filling placeholder %s for key %s: %s", ph_idx, key, e
                )

        # --- HEADER placeholders (styled) ---
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text  # type: ignore

            if "NOM_PROJET" in text:
                shape.text = safe_get("project_name")  # type: ignore
                for p in shape.text_frame.paragraphs:  # type: ignore
                    for r in p.runs:
                        r.font.size = Pt(24)
                        r.font.bold = True
                        r.font.color.rgb = RGBColor(255, 255, 255)
            elif "personnes" in text:
                shape.text = safe_get("num_people", " personnes")  # type: ignore
            elif "Mois_debut" in text or "Mois_fin" in text:
                start = safe_get("start_date")
                end = safe_get("end_date")
                shape.text = f"{start} - {end}"  # type: ignore
            elif "Enjeux" in text or "€" in text:
                shape.text = (  # type: ignore
                    f"Enjeux financier : {safe_get('budget_k_eur')}k€"  # type: ignore
                )
            else:
                continue

            # Non-title header fields: 16pt white, non-bold
            if "NOM_PROJET" not in text:
                for p in shape.text_frame.paragraphs:  # type: ignore
                    for r in p.runs:
                        r.font.size = Pt(16)
                        r.font.bold = False
                        r.font.color.rgb = RGBColor(255, 255, 255)

        prs.save(str(output_path))
        logger.info("PowerPoint generated: %s", output_path)
        return output_path

    # -----------------------------
    # Reasoning / generation step
    # -----------------------------
    async def _run_reasoning_step(self, state: MessagesState):
        if self.model is None:
            raise RuntimeError("Model not initialized. Call async_init().")
        # Check for prompt initialization
        if self.prompt_template is None or self.summary_prompt_template is None:
            raise RuntimeError("Prompts not initialized. Call async_init().")

        last = state["messages"][-1]
        if not isinstance(last.content, str):
            raise TypeError("Expected string content for the last message.")
        question = last.content

        try:
            doc_tag_ids = get_document_library_tags_ids(self.get_runtime_context())
            document_uids = get_document_uids(self.get_runtime_context())
            search_policy = get_search_policy(self.get_runtime_context())
            top_k = self.get_tuned_int("rag.top_k", default=6)

            runtime_ctx = self.get_runtime_context()
            include_session_scope, include_corpus_scope = get_vector_search_scopes(
                runtime_ctx
            )
            if not runtime_ctx or not runtime_ctx.session_id:
                raise RuntimeError(
                    "Runtime context missing session_id; required for scoped retrieval."
                )
            hits: List[VectorSearchHit] = await self.search_client.search(
                question=question,
                top_k=top_k,
                document_library_tags_ids=doc_tag_ids,
                document_uids=document_uids,
                search_policy=search_policy,
                session_id=runtime_ctx.session_id,
                include_session_scope=include_session_scope,
                include_corpus_scope=include_corpus_scope,
            )

            if not hits:
                warn = "I did not find any relevant documents. Try rephrasing your question."
                messages = await self.with_chat_context_text(
                    [HumanMessage(content=warn)]
                )
                return {"messages": [await self.model.ainvoke(messages)]}

            hits = sort_hits(hits)
            ensure_ranks(hits)

            sys_msg = SystemMessage(content=self._system_prompt())
            context_text = format_sources_for_prompt(hits, snippet_chars=1200)

            # --- Use Tuned Extraction Prompt ---
            human_msg = HumanMessage(
                content=self.prompt_template.format_messages(
                    format_instructions=self.format_instructions,
                    context=context_text,
                )[0].content
            )
            # ----------------------------------

            messages = [sys_msg, human_msg]
            messages = await self.with_chat_context_text(messages)

            answer_msg = await self.model.ainvoke(messages)
            answer_text = getattr(answer_msg, "content", "")
            try:
                assert self.parser is not None
                parsed = self.parser.parse(answer_text)
                # JsonOutputParser with pydantic_object returns a BaseModel instance.
                # Convert it to a plain dict for downstream `.get` and mapping usage.
                if hasattr(parsed, "model_dump"):
                    structured_data = parsed.model_dump()
                elif hasattr(parsed, "dict"):
                    structured_data = parsed.dict()
                else:
                    # Fallback: keep as-is (may already be a dict-like)
                    structured_data = parsed  # type: ignore[assignment]
            except Exception as e:
                logger.warning("Structured parse failed: %s", e)
                structured_data = {"raw_text": answer_text}

            attach_sources_to_llm_response(answer_msg, hits)

            project_name = structured_data.get("project_name", "project_auto").replace(
                " ", "_"
            )
            ppt_path: Optional[Path] = None
            pdf_path: Optional[Path] = None
            pdf_download_url: Optional[str] = None

            try:
                ppt_path = await self._fill_ppt_template(structured_data, project_name)
            except Exception as e:
                logger.exception("Failed to generate PPTX: %s", e)
                error_msg = AIMessage(content=f"❌ Error generating PowerPoint: {e}")
                return {
                    "messages": [answer_msg, error_msg],
                    "structured_data": structured_data,
                }

            # Upload PPTX
            try:
                user_id = self.get_end_user_id()
                timestamp = int(time.time() * 1000)
                unique_ppt_name = f"fiche_ref_{project_name}_{timestamp}.pptx"
                unique_pdf_name = f"fiche_ref_{project_name}_{timestamp}.pdf"
                with open(ppt_path, "rb") as f:
                    upload_result = await self.upload_user_blob(
                        key=f"{user_id}_{unique_ppt_name}",
                        file_content=f,
                        filename=unique_ppt_name,
                        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                download_url = upload_result.download_url

                # --- Convert to PDF and upload ---
                pdf_path = self._convert_pptx_to_pdf(ppt_path)
                if pdf_path and pdf_path.exists():
                    with open(pdf_path, "rb") as f:
                        pdf_upload = await self.upload_user_blob(
                            key=f"{user_id}_{unique_pdf_name}",
                            file_content=f,
                            filename=unique_pdf_name,
                            content_type="application/pdf",
                        )
                        logger.info(f"pdf_upload: {pdf_upload}")
                    pdf_download_url = pdf_upload.download_url

                # --- Use Tuned Summary Prompt ---
                summary_template_text = self.summary_prompt_template.format_messages(
                    structured_data=structured_data
                )[0].content
                summary_prompt = HumanMessage(content=summary_template_text)
                # ------------------------------

                summary_msg = await self.model.ainvoke([summary_prompt])
                summary_text = getattr(summary_msg, "content", "").strip()

                text_summary = "✅ **Slides generated successfully!**"
                if summary_text:
                    text_summary += (
                        "\n\n🔎 **Generated content summary:**\n\n" + summary_text
                    )

                parts: List[MessagePart] = [
                    TextPart(text=text_summary),
                    LinkPart(
                        href=download_url,
                        title=f"Download {upload_result.file_name}",
                        kind=LinkKind.download,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    ),
                ]

                if pdf_download_url:
                    parts.append(
                        LinkPart(
                            title="View (PDF Preview)",
                            kind=LinkKind.view,
                            mime="application/pdf",
                            document_uid=pdf_upload.document_uid,
                            file_name=pdf_upload.file_name,
                        )
                    )

                ai_msg = AIMessage(content="", parts=parts)
                return {"messages": [ai_msg], "structured_data": structured_data}

            except WorkspaceRetrievalError as e:
                logger.exception("Asset upload error: %s", e)
                return {
                    "messages": [AIMessage(content=f"❌ Upload error: {e}")],
                    "structured_data": structured_data,
                }

            finally:
                for p in [ppt_path, pdf_path]:
                    if p and p.exists():
                        try:
                            p.unlink(missing_ok=True)
                        except Exception as e:
                            logger.warning(
                                "Failed to delete temporary file %s: %s", p, e
                            )

        except Exception:
            logger.exception("Sloan: error in reasoning step.")
            fallback = await self.model.ainvoke(
                [
                    HumanMessage(
                        content="An error occurred while extracting the information."
                    )
                ]
            )
            return {"messages": [fallback]}
