# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, List, Tuple

from pptx import Presentation

from knowledge_flow_backend.core.processors.input.common.base_input_processor import BaseMarkdownProcessor, InputConversionError
from knowledge_flow_backend.core.processors.input.lightweight_markdown_processor.base_lite_md_processor import BaseLiteMdProcessor
from knowledge_flow_backend.core.processors.input.lightweight_markdown_processor.lite_markdown_structures import (
    LiteMarkdownOptions,
    LiteMarkdownResult,
    LitePageMarkdown,
    collapse_whitespace,
    enforce_max_chars,
)

logger = logging.getLogger(__name__)


def _escape_pipes(s: str) -> str:
    return s.replace("|", r"\|")


def _clean_text(s: str) -> str:
    # Keep output stable for downstream markdown rendering.
    # Replace control chars with spaces (except common whitespace).
    if not s:
        return ""
    return "".join(ch if (ch >= " " or ch in "\n\t") else " " for ch in s)


class LitePptxToMdExtractor(BaseLiteMdProcessor):
    """
    Lightweight PPTX -> Markdown (python-pptx only).

    Guarantees:
    - no markitdown
    - no images
    - no external binaries
    - deterministic slide-wise extraction
    """

    description = "Fast PPTX-to-Markdown extractor (python-pptx only)."

    def extract(self, file_path: Path, options: LiteMarkdownOptions | None = None) -> LiteMarkdownResult:
        opts = options or LiteMarkdownOptions()

        pres = Presentation(str(file_path))
        slide_count = len(pres.slides)

        start_s, end_s = self._safe_slide_range(slide_count, opts.page_range)
        pages: List[LitePageMarkdown] = []

        for s_idx in range(start_s, end_s + 1):
            slide = pres.slides[s_idx - 1]
            lines: List[str] = []

            if opts.add_page_headings:
                title_text = ""
                try:
                    if slide.shapes.title and slide.shapes.title.text:
                        title_text = slide.shapes.title.text.strip()
                except Exception:
                    title_text = ""
                heading = f"## Slide {s_idx}" + (f": {title_text}" if title_text else "")
                lines.append(heading)
                lines.append("")

            for shape in slide.shapes:
                try:
                    part = self._shape_to_markdown(shape, opts)
                    if part:
                        lines.append(part)
                except Exception as e:  # noqa: BLE001
                    logger.debug(
                        "[LITE_PPTX][IMPLEM] shape export skipped on slide %s: %s",
                        s_idx,
                        e,
                    )

            slide_md = "\n".join(line for line in lines if line is not None)
            slide_md = _clean_text(slide_md)

            if opts.normalize_whitespace:
                slide_md = collapse_whitespace(slide_md)

            slide_md = slide_md.strip()
            pages.append(LitePageMarkdown(page_no=s_idx, markdown=slide_md, char_count=len(slide_md)))

        combined = "\n\n".join(p.markdown for p in pages)
        combined, truncated = enforce_max_chars(combined, opts.max_chars)

        if truncated and opts.return_per_page:
            kept: List[LitePageMarkdown] = []
            acc = 0
            for i, p in enumerate(pages):
                sep = 2 if i > 0 else 0
                if acc + sep + len(p.markdown) > len(combined):
                    break
                acc += sep + len(p.markdown)
                kept.append(p)
            pages = kept

        return LiteMarkdownResult(
            document_name=file_path.name,
            page_count=slide_count,
            total_chars=len(combined),
            truncated=truncated,
            markdown=combined,
            pages=pages if opts.return_per_page else [],
            extras={"engine": "python-pptx-slidewise"},
        )

    def _safe_slide_range(self, slide_count: int, page_range: Tuple[int, int] | None) -> Tuple[int, int]:
        if not page_range:
            return (1, slide_count)
        start_s, end_s = page_range
        start_s = max(1, min(int(start_s), slide_count))
        end_s = max(start_s, min(int(end_s), slide_count))
        return (start_s, end_s)

    def _shape_to_markdown(self, shape, opts: LiteMarkdownOptions) -> str:
        # Tables
        has_table = bool(getattr(shape, "has_table", False))
        table = getattr(shape, "table", None)
        if has_table and table is not None:
            return self._table_to_markdown(table, opts)

        # Text frames
        has_text_frame = bool(getattr(shape, "has_text_frame", False))
        text_frame = getattr(shape, "text_frame", None)
        if has_text_frame and text_frame is not None:
            return self._text_frame_to_markdown(text_frame)

        # Some shapes expose plain .text
        text_value = getattr(shape, "text", None)
        if text_value:
            return str(text_value).strip()

        return ""

    def _text_frame_to_markdown(self, text_frame) -> str:
        out: List[str] = []
        for para in text_frame.paragraphs:
            text = (getattr(para, "text", "") or "").strip()
            if not text:
                continue

            level = int(getattr(para, "level", 0) or 0)
            indent = "  " * max(0, level)

            # Robust, simple rule: bullet only when level > 0
            if level > 0:
                out.append(f"{indent}- {text}")
            else:
                out.append(f"{indent}{text}")

        return "\n".join(out).strip()

    def _table_to_markdown(self, table, opts: LiteMarkdownOptions) -> str:
        try:
            rows = len(table.rows)
            cols = len(table.columns)
        except Exception:
            return ""

        if rows <= 0 or cols <= 0:
            return ""

        max_r = min(rows, max(0, opts.max_table_rows) + 1)  # keep header row
        max_c = min(cols, max(0, opts.max_table_cols))
        if max_c <= 0:
            return ""

        def cell_text(r: int, c: int) -> str:
            try:
                cell = table.cell(r, c)
                raw = ""
                try:
                    raw = cell.text or ""
                except Exception:
                    raw = ""
                raw = raw.replace("\r", " ").replace("\n", " ")
                raw = " ".join(raw.split())
                return _escape_pipes(raw.strip())
            except Exception:
                return ""

        lines: List[str] = []
        lines.append("")

        header = [cell_text(0, c) for c in range(max_c)]
        header = [h if h else " " for h in header]
        lines.append("| " + " | ".join(header) + " |")
        lines.append("| " + " | ".join(["---"] * max_c) + " |")

        for r in range(1, max_r):
            row_cells = [cell_text(r, c) for c in range(max_c)]
            row_cells = [v if v else " " for v in row_cells]
            lines.append("| " + " | ".join(row_cells) + " |")

        if max_r < rows or max_c < cols:
            lines.append("… (table truncated)")

        lines.append("")
        return "\n".join(lines).strip()


class LitePptxMarkdownProcessor(BaseMarkdownProcessor):
    """
    Adapter so the lightweight PPTX extractor can be used as a BaseMarkdownProcessor.
    """

    description = "Lightweight PPTX ingestion that converts slides to Markdown previews."

    def __init__(self) -> None:
        super().__init__()
        self._lite = LitePptxToMdExtractor()

    def check_file_validity(self, file_path: Path) -> bool:
        try:
            pres = Presentation(str(file_path))
            if len(pres.slides) == 0:
                logger.warning("LitePptxMarkdownProcessor: PPTX %s has no slides.", file_path)
                return False
            return True
        except Exception as e:  # noqa: BLE001
            logger.error("LitePptxMarkdownProcessor: invalid PPTX %s: %s", file_path, e)
            return False

    def extract_file_metadata(self, file_path: Path) -> dict[str, Any]:
        try:
            pres = Presentation(str(file_path))
            slide_count = len(pres.slides)
            metadata: dict[str, Any] = {
                "document_name": file_path.name,
                "title": file_path.stem,
                "page_count": slide_count,
            }
            return metadata
        except Exception as e:  # noqa: BLE001
            logger.error("LitePptxMarkdownProcessor: error extracting metadata from %s: %s", file_path, e)
            return {"document_name": file_path.name, "error": str(e)}

    def convert_file_to_markdown(self, file_path: Path, output_dir: Path, document_uid: str | None) -> dict:
        output_markdown_path = output_dir / "output.md"
        try:
            result = self._lite.extract(file_path, LiteMarkdownOptions())
            markdown = result.markdown or ""
            output_dir.mkdir(parents=True, exist_ok=True)
            output_markdown_path.write_text(markdown, encoding="utf-8")

            engine = (result.extras or {}).get("engine") if isinstance(result.extras, dict) else None
            message = f"Lite PPTX conversion succeeded (engine={engine})" if engine else "Lite PPTX conversion succeeded."
        except Exception as exc:  # noqa: BLE001
            logger.exception("LitePptxMarkdownProcessor: conversion failed for %s", file_path)
            raise InputConversionError(f"LitePptxMarkdownProcessor failed for '{file_path.name}': {exc}") from exc

        return {
            "doc_dir": str(output_dir),
            "md_file": str(output_markdown_path),
            "message": message,
        }
