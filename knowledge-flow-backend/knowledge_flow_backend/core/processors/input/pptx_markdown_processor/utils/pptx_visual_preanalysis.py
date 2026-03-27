# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Pre-analyzes PPTX slides to decide whether vision enrichment is needed."""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Any, Optional

from pptx.enum.shapes import MSO_SHAPE_TYPE

from knowledge_flow_backend.core.processors.input.pptx_markdown_processor.utils.pptx_native_slide_extractor import (
    NativeSlideContent,
)


@dataclass
class PptxVisualElementSummary:
    shape_type: str
    shape_type_raw: Optional[str] = None
    shape_name: Optional[str] = None
    has_text: bool = False
    text_preview: Optional[str] = None
    is_placeholder: bool = False
    placeholder_type: Optional[str] = None
    top: Optional[int] = None
    left: Optional[int] = None
    width: Optional[int] = None
    height: Optional[int] = None
    area: Optional[int] = None


@dataclass
class PptxSlidePreanalysis:
    slide_number: int
    visual_elements: list[PptxVisualElementSummary] = field(default_factory=list)
    text_shape_count: int = 0
    table_count: int = 0
    picture_count: int = 0
    group_count: int = 0
    chart_count: int = 0
    other_count: int = 0
    raw_text_length: int = 0
    bullet_count: int = 0
    visual_area_ratio: float = 0.0
    needs_vision: bool = False
    vision_reasons: list[str] = field(default_factory=list)
    vision_priority: str = "low"


@dataclass
class PptxDocumentPreanalysis:
    total_slides: int
    slides: list[PptxSlidePreanalysis] = field(default_factory=list)
    activate_vision: bool = False
    slides_for_vision: list[int] = field(default_factory=list)


def _clean_text(text: str) -> str:
    if not text:
        return ""
    return " ".join(text.replace("\r", " ").replace("\n", " ").split()).strip()


def _extract_shape_text_lines(shape: Any) -> list[str]:
    has_text_frame = bool(getattr(shape, "has_text_frame", False))
    text_frame = getattr(shape, "text_frame", None)
    if has_text_frame and text_frame is not None:
        lines: list[str] = []
        for para in getattr(text_frame, "paragraphs", []):
            text = _clean_text(getattr(para, "text", "") or "")
            if text:
                lines.append(text)
        return lines

    text_value = _clean_text(getattr(shape, "text", "") or "")
    return [text_value] if text_value else []


def _classify_shape(shape: Any) -> str:
    """Classify a PPTX shape for vision-oriented pre-analysis."""
    if bool(getattr(shape, "has_table", False)):
        return "table"

    if bool(getattr(shape, "has_text_frame", False)):
        return "text"

    shape_type = getattr(shape, "shape_type", None)
    if shape_type is None:
        return "other"

    if shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "picture"
    if shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if shape_type == MSO_SHAPE_TYPE.CHART:
        return "chart"

    return "other"


def _build_visual_element_summary(shape: Any, shape_type: str) -> PptxVisualElementSummary:
    text_lines = _extract_shape_text_lines(shape)
    text_preview = " ".join(text_lines[:2]).strip() if text_lines else None
    if text_preview and len(text_preview) > 120:
        text_preview = text_preview[:117] + "..."

    top = getattr(shape, "top", None)
    left = getattr(shape, "left", None)
    width = getattr(shape, "width", None)
    height = getattr(shape, "height", None)

    area = None
    if width is not None and height is not None:
        try:
            area = int(width) * int(height)
        except Exception:
            area = None

    is_placeholder = bool(getattr(shape, "is_placeholder", False))
    placeholder_type = None
    if is_placeholder:
        try:
            placeholder_type = str(getattr(shape.placeholder_format, "type", None))
        except Exception:
            placeholder_type = None

    shape_type_raw = None
    try:
        shape_type_raw = str(getattr(shape, "shape_type", None))
    except Exception:
        shape_type_raw = None

    return PptxVisualElementSummary(
        shape_type=shape_type,
        shape_type_raw=shape_type_raw,
        shape_name=getattr(shape, "name", None),
        has_text=bool(text_lines),
        text_preview=text_preview,
        is_placeholder=is_placeholder,
        placeholder_type=placeholder_type,
        top=top,
        left=left,
        width=width,
        height=height,
        area=area,
    )


def _compute_visual_area_ratio(slide: Any, visual_elements: list[PptxVisualElementSummary]) -> float:
    try:
        presentation = slide.part.package.presentation_part.presentation
        slide_width = int(getattr(presentation, "slide_width", 0) or 0)
        slide_height = int(getattr(presentation, "slide_height", 0) or 0)
    except Exception:
        slide_width = 0
        slide_height = 0

    slide_area = slide_width * slide_height
    if slide_area <= 0:
        return 0.0

    visual_area = sum(element.area or 0 for element in visual_elements)
    return min(1.0, visual_area / slide_area)


def _compute_vision_priority(slide_summary: PptxSlidePreanalysis) -> str:
    """
    Estimate whether a slide should be sent to vision in the current V1 setup.

    - high: vision is likely to add meaningful information beyond native extraction
    - low: the slide is visual, but vision is unlikely to add enough value for now
    """
    reasons = set(slide_summary.vision_reasons)

    if "chart" in reasons:
        return "high"

    if "visual_heavy" in reasons and ("object_placeholder" in reasons or "unclassified_visual" in reasons or "group" in reasons or "picture" in reasons):
        return "high"

    if "low_native_text_with_visuals" in reasons and len(slide_summary.visual_elements) >= 1:
        return "high"

    if "group" in reasons:
        return "high"

    if "object_placeholder" in reasons or "unclassified_visual" in reasons:
        return "high"

    return "low"


def _compute_vision_reasons(slide_summary: PptxSlidePreanalysis) -> list[str]:
    reasons: list[str] = []

    if slide_summary.picture_count > 0:
        reasons.append("picture")
    if slide_summary.chart_count > 0:
        reasons.append("chart")
    if slide_summary.group_count > 0:
        reasons.append("group")
    if any(element.placeholder_type and "OBJECT" in element.placeholder_type for element in slide_summary.visual_elements):
        reasons.append("object_placeholder")
    if slide_summary.other_count > 0 and slide_summary.picture_count == 0 and slide_summary.chart_count == 0:
        reasons.append("unclassified_visual")
    # These flags do not describe the visual object type itself.
    # They explain why visual enrichment is likely to add value.
    if slide_summary.visual_area_ratio > 0.25:
        reasons.append("visual_heavy")
    if slide_summary.raw_text_length < 120 and len(slide_summary.visual_elements) >= 2:
        reasons.append("low_native_text_with_visuals")

    return reasons


def preanalyze_slide(
    slide: Any,
    slide_number: int,
    native_content: Optional[NativeSlideContent] = None,
) -> PptxSlidePreanalysis:
    summary = PptxSlidePreanalysis(slide_number=slide_number)

    if native_content is not None:
        summary.bullet_count = len(native_content.bullets)
        summary.raw_text_length = sum(len(block) for block in native_content.raw_text_blocks)
        summary.table_count = len(native_content.tables)

    for shape in getattr(slide, "shapes", []):
        shape_kind = _classify_shape(shape)

        if shape_kind == "text":
            summary.text_shape_count += 1
            continue

        if shape_kind == "table":
            # Already handled natively; no need to route through vision by default.
            continue

        element = _build_visual_element_summary(shape, shape_kind)
        summary.visual_elements.append(element)

        if shape_kind == "picture":
            summary.picture_count += 1
        elif shape_kind == "group":
            summary.group_count += 1
        elif shape_kind == "chart":
            summary.chart_count += 1
        else:
            summary.other_count += 1

    summary.visual_area_ratio = _compute_visual_area_ratio(slide, summary.visual_elements)
    summary.vision_reasons = _compute_vision_reasons(summary)
    summary.needs_vision = bool(summary.vision_reasons)
    summary.vision_priority = _compute_vision_priority(summary)

    return summary


def preanalyze_presentation(
    presentation: Any,
    native_contents: Optional[list[NativeSlideContent]] = None,
) -> PptxDocumentPreanalysis:
    slides = list(getattr(presentation, "slides", []))
    result = PptxDocumentPreanalysis(total_slides=len(slides))

    native_by_number = {}
    if native_contents:
        native_by_number = {content.slide_number: content for content in native_contents}

    for slide_number, slide in enumerate(slides, start=1):
        slide_summary = preanalyze_slide(
            slide,
            slide_number,
            native_content=native_by_number.get(slide_number),
        )
        result.slides.append(slide_summary)

    result.slides_for_vision = [slide_summary.slide_number for slide_summary in result.slides if slide_summary.needs_vision]
    result.activate_vision = bool(result.slides_for_vision)

    return result
