# Copyright Thales 2025
#
# Licensed under the Apache License, Version 2.0 (the "License");
# you may not use this file except in compliance with the License.
# You may obtain a copy of the License at
#
#     http://www.apache.org/licenses/LICENSE-2.0
#
# Unless required by applicable law or agreed to in writing, software
# distributed under the License is distributed on an "AS IS" BASIS,
# WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
# See the License for the specific language governing permissions and
# limitations under the License.

"""Main PPTX Markdown processor.

This processor orchestrates PPTX-specific helper modules for native slide extraction,
formatting, speaker notes, and deck-level cleanup. The split keeps the native
extraction reusable for future vision-enriched PPTX processing.
"""

import logging
from pathlib import Path
from typing import Any

from pptx import Presentation

from knowledge_flow_backend.application_context import get_configuration
from knowledge_flow_backend.common.processing_profile_context import get_current_processing_profile
from knowledge_flow_backend.common.structures import IngestionProcessingProfile
from knowledge_flow_backend.core.processors.input.common.base_input_processor import BaseMarkdownProcessor, InputConversionError
from knowledge_flow_backend.core.processors.input.common.image_describer import (
    PPTX_MEDIUM_VISION_DESCRIBE_PROMPT_V1,
    build_image_describer,
)
from knowledge_flow_backend.core.processors.input.pptx_markdown_processor.utils.pptx_deck_noise import (
    detect_repeated_noise_texts,
)
from knowledge_flow_backend.core.processors.input.pptx_markdown_processor.utils.pptx_native_slide_extractor import (
    extract_native_slide_content,
)
from knowledge_flow_backend.core.processors.input.pptx_markdown_processor.utils.pptx_slide_markdown_formatter import (
    format_slide_markdown,
)
from knowledge_flow_backend.core.processors.input.pptx_markdown_processor.utils.pptx_vision_enricher import (
    enrich_slides_with_vision,
)
from knowledge_flow_backend.core.processors.input.pptx_markdown_processor.utils.pptx_visual_preanalysis import (
    preanalyze_presentation,
)

logger = logging.getLogger(__name__)


class PptxMarkdownProcessor(BaseMarkdownProcessor):
    description = "Converts PPTX slide decks into Markdown sections, slide by slide."

    def __init__(self):
        super().__init__()
        self.image_describer = None
        self._warned_missing_vision_model = False

    def _resolve_effective_options(self) -> tuple[IngestionProcessingProfile, bool]:
        processing = get_configuration().processing
        current_profile = get_current_processing_profile()
        active_profile = processing.normalize_profile(current_profile)

        # V1 choice:
        # - fast: native only
        # - medium: enable vision enrichment
        # - rich: kept disabled for now until the rich strategy is finalized
        enable_vision = active_profile == IngestionProcessingProfile.MEDIUM

        return active_profile, enable_vision

    def _resolve_image_describer(self, enable_vision: bool):
        if not enable_vision:
            return None
        if self.image_describer is not None:
            return self.image_describer
        if not get_configuration().vision_model:
            if not self._warned_missing_vision_model:
                logger.warning("[PROCESSOR][PPTX] Vision model configuration is missing while vision enrichment is enabled.")
                self._warned_missing_vision_model = True
            return None
        self.image_describer = build_image_describer(
            get_configuration().vision_model,
            system_prompt=PPTX_MEDIUM_VISION_DESCRIBE_PROMPT_V1,
        )
        return self.image_describer

    def _build_native_markdown(
        self,
        presentation: Any,
        visual_enrichments: dict[int, str] | None = None,
    ) -> tuple[list[str], list]:
        slide_markdowns = []
        native_contents = []

        visual_enrichments = visual_enrichments or {}
        repeated_noise_texts = detect_repeated_noise_texts(presentation.slides)

        for slide_number, slide in enumerate(presentation.slides, start=1):
            native_content = extract_native_slide_content(
                slide,
                slide_number,
                repeated_noise_texts=repeated_noise_texts,
            )
            native_contents.append(native_content)

            slide_md = format_slide_markdown(
                native_content,
                visual_enrichment=visual_enrichments.get(slide_number),
            )

            if slide_md:
                slide_markdowns.append(slide_md)

        return slide_markdowns, native_contents

    def _log_visual_preanalysis(self, visual_preanalysis) -> None:
        logger.info(
            "[PPTX][PREANALYSIS] activate_vision=%s slides_for_vision=%s",
            visual_preanalysis.activate_vision,
            visual_preanalysis.slides_for_vision,
        )
        for slide_summary in visual_preanalysis.slides:
            if not slide_summary.needs_vision:
                continue

            logger.info(
                "[PPTX][PREANALYSIS] slide=%s needs_vision=%s vision_priority=%s reasons=%s visual_area_ratio=%.3f pictures=%s groups=%s charts=%s others=%s",
                slide_summary.slide_number,
                slide_summary.needs_vision,
                slide_summary.vision_priority,
                slide_summary.vision_reasons,
                slide_summary.visual_area_ratio,
                slide_summary.picture_count,
                slide_summary.group_count,
                slide_summary.chart_count,
                slide_summary.other_count,
            )

    def _select_slides_for_vision(self, visual_preanalysis) -> list[int]:
        return [slide_summary.slide_number for slide_summary in visual_preanalysis.slides if slide_summary.needs_vision]

    def check_file_validity(self, file_path: Path) -> bool:
        """Checks if the PPTX file is valid and can be opened."""
        try:
            Presentation(str(file_path))
            return True
        except Exception as e:
            logger.error(f"Invalid or corrupted PPTX file: {file_path} - {e}")
            return False

    def extract_file_metadata(self, file_path: Path) -> dict[str, Any]:
        """Extracts basic metadata from the PPTX file."""
        metadata: dict[str, Any] = {"document_name": file_path.name}
        try:
            presentation = Presentation(str(file_path))
            metadata["num_slides"] = len(presentation.slides)
        except Exception as e:
            logger.error(f"Error reading PPTX file: {e}")
            metadata["error"] = str(e)
        return metadata

    def convert_file_to_markdown(self, file_path: Path, output_dir: Path, document_uid: str | None) -> dict:
        """Converts each slide's content into structured Markdown."""
        output_dir.mkdir(parents=True, exist_ok=True)
        md_path = output_dir / "output.md"

        try:
            active_profile, enable_vision = self._resolve_effective_options()
            image_describer = self._resolve_image_describer(enable_vision)

            logger.info(
                "[PROCESSOR][PPTX] Using profile=%s enable_vision=%s vision_model_available=%s",
                active_profile.value,
                enable_vision,
                bool(image_describer),
            )

            presentation = Presentation(str(file_path))

            _, native_contents = self._build_native_markdown(presentation)

            visual_preanalysis = preanalyze_presentation(
                presentation,
                native_contents=native_contents,
            )
            self._log_visual_preanalysis(visual_preanalysis)

            visual_enrichments: dict[int, str] = {}
            if enable_vision and image_describer:
                slides_to_enrich = self._select_slides_for_vision(visual_preanalysis)
                logger.info(
                    "[PROCESSOR][PPTX] Selected %s slide(s) for vision enrichment: %s",
                    len(slides_to_enrich),
                    slides_to_enrich,
                )

                visual_enrichments = enrich_slides_with_vision(
                    pptx_path=file_path,
                    slide_numbers=slides_to_enrich,
                    output_dir=output_dir,
                    image_describer=image_describer,
                )

            slide_markdowns = [
                format_slide_markdown(
                    native_content,
                    visual_enrichment=visual_enrichments.get(native_content.slide_number),
                )
                for native_content in native_contents
            ]

            content = "\n\n---\n\n".join(slide_markdowns) if slide_markdowns else "*No extractable text*"
            md_path.write_text(content, encoding="utf-8")

            return {
                "doc_dir": str(output_dir),
                "md_file": str(md_path),
                "message": "PPTX slides converted to structured Markdown.",
            }

        except Exception as exc:
            logger.exception("Failed to convert PPTX to Markdown: %s", file_path)
            raise InputConversionError(f"PptxMarkdownProcessor failed for '{file_path.name}': {exc}") from exc
